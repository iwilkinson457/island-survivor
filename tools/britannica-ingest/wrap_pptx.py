"""
PPTX -> Markdown wrapper for Britannica ingestion.

Uses python-pptx (already installed).
Extracts slide titles and text content. Notes are included if present.

Limitations:
- Images, charts, SmartArt, and embedded media are ignored.
- Complex animations and transitions are not represented.
- Heavily image-based decks will yield very little text (expect grade C/D).
"""

from pathlib import Path


def extract(source: Path) -> dict:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return {"error": "python-pptx not installed (pip install python-pptx)"}

    try:
        prs = Presentation(str(source))
    except Exception as e:
        return {"error": f"python-pptx failed to open file: {e}"}

    slides_md = []
    all_text_parts = []
    deck_title = None

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_title = None
        text_blocks = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Detect title placeholder
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    continue
                try:
                    ph = shape.placeholder_format
                    is_title = ph is not None and ph.idx in (0, 1)
                except Exception:
                    is_title = False

                if is_title and not slide_title:
                    slide_title = text
                    if slide_num == 1 and not deck_title:
                        deck_title = text
                    text_blocks.append(f"**{text}**")
                else:
                    text_blocks.append(f"- {text}")
                all_text_parts.append(text)

        # Include speaker notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                notes = notes_tf.text.strip()
                if notes:
                    notes_text = f"\n  > _Notes: {notes[:300]}_"

        if text_blocks or notes_text:
            heading = slide_title or f"Slide {slide_num}"
            slides_md.append(f"### Slide {slide_num}: {heading}\n\n" + "\n".join(text_blocks) + notes_text)

    body_md = "\n\n".join(slides_md)
    full_text = " ".join(all_text_parts)

    # Summary
    substantive = [t for t in all_text_parts if len(t) > 25]
    summary = " ".join(substantive[:4])[:500].strip()
    if not summary:
        summary = full_text[:300].strip()

    # Key facts
    key_facts = _extract_key_facts(all_text_parts, len(prs.slides))

    # Operational relevance
    op_relevance = _detect_operational_relevance(full_text)

    # Gaps
    gaps = []
    slide_count = len(prs.slides)
    if not all_text_parts:
        gaps.append("No text extracted — deck may be entirely image-based")
    elif len(all_text_parts) < 5:
        gaps.append("Very little text extracted — deck is likely mostly images")
    if slide_count > 30:
        gaps.append(f"Large deck ({slide_count} slides) — content may be incomplete at limits")

    return {
        "title": deck_title or source.stem,
        "summary": summary,
        "key_facts": key_facts,
        "operational_relevance": op_relevance,
        "gaps": gaps,
        "method": f"python-pptx ({slide_count} slides)",
        "body_md": body_md,
    }


def _extract_key_facts(text_parts: list, slide_count: int) -> list:
    facts = [f"Presentation has {slide_count} slide(s)"]
    seen = set()
    for text in text_parts:
        text = text.strip()
        if len(text) < 15 or text in seen:
            continue
        seen.add(text)
        facts.append(text[:200])
        if len(facts) >= 9:
            break
    return facts


def _detect_operational_relevance(body: str) -> str:
    lower = body.lower()
    signals = {
        "procedure": "Contains procedural content",
        "maintenance": "Relates to maintenance",
        "safety": "Contains safety information",
        "training": "Training or instructional content",
        "project": "Project-related content",
        "overview": "Overview or briefing deck",
        "report": "Report or status deck",
        "specification": "Technical specification",
    }
    matched = [v for k, v in signals.items() if k in lower]
    return "; ".join(matched[:3]) if matched else "Presentation — review slides for specific relevance"
